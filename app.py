from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify, send_file
from werkzeug.security import generate_password_hash, check_password_hash
from flask_sqlalchemy import SQLAlchemy
from functools import wraps
import os
import io
import json
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from ai_service import generate_content_section, refine_content_section, generate_outline
import time

app = Flask(__name__)
app.config['SECRET_KEY'] = os.urandom(24).hex()

# PostgreSQL Database Configuration
DATABASE_URL = os.getenv('DATABASE_URL')
if not DATABASE_URL:
    # Using provided credentials: database=docxbuilder, password=PostgreSQL1036
    # Default: username=postgres, host=localhost, port=5432
    # Using psycopg3 driver (postgresql+psycopg://) for psycopg package
    DATABASE_URL = 'postgresql+psycopg://postgres:PostgreSQL1036@localhost:5432/docxbuilder'

app.config['SQLALCHEMY_DATABASE_URI'] = DATABASE_URL
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)


# User Model
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(255), nullable=False)
    
    # Relationships
    revisions = db.relationship('Revision', backref='user', lazy=True)
    comments = db.relationship('Comment', backref='user', lazy=True)
    feedback = db.relationship('Feedback', backref='user', lazy=True)

    def __repr__(self):
        return f'<User {self.username}>'


# Content Model for storing pasted content
class Content(db.Model):
    password = db.Column(db.String(255), primary_key=True)  # Password as Primary Key
    username = db.Column(db.String(80), nullable=False)
    content_entered = db.Column(db.Text, nullable=False)  # Content pasted by user
    content_number = db.Column(db.Integer, nullable=False, unique=True)  # Content number (auto-incremented manually)
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)

    def __repr__(self):
        return f'<Content #{self.content_number} by {self.username}>'


# Project Model for AI-powered document generation
class Project(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    project_type = db.Column(db.String(20), nullable=False)  # 'word' or 'powerpoint'
    main_topic = db.Column(db.Text, nullable=False)
    username = db.Column(db.String(80), nullable=False)
    outline = db.Column(db.Text)  # JSON string of sections/slides
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow, nullable=False)

    # Relationships
    sections = db.relationship('Section', backref='project', lazy=True, cascade='all, delete-orphan')
    revisions = db.relationship('Revision', backref='project', lazy=True, cascade='all, delete-orphan')
    comments = db.relationship('Comment', backref='project', lazy=True, cascade='all, delete-orphan')

    def __repr__(self):
        return f'<Project {self.name} ({self.project_type})>'


# Section Model for document sections/slides
class Section(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    project_id = db.Column(db.Integer, db.ForeignKey('project.id'), nullable=False)
    section_number = db.Column(db.Integer, nullable=False)
    title = db.Column(db.String(500), nullable=False)
    content = db.Column(db.Text)  # AI-generated or user-edited content
    section_type = db.Column(db.String(20), nullable=False)  # 'section' for Word, 'slide' for PowerPoint
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow, nullable=False)

    # Relationships
    revisions = db.relationship('Revision', backref='section', lazy=True, cascade='all, delete-orphan', order_by='Revision.created_at.desc()')
    feedback = db.relationship('Feedback', backref='section', lazy=True, cascade='all, delete-orphan', uselist=False)
    comments = db.relationship('Comment', backref='section', lazy=True, cascade='all, delete-orphan', order_by='Comment.created_at.desc()')

    def __repr__(self):
        return f'<Section {self.section_number}: {self.title}>'


# Revision Model for tracking content changes
class Revision(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=False)
    project_id = db.Column(db.Integer, db.ForeignKey('project.id'), nullable=True)  # Made nullable temporarily for migration
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)  # User who made the revision
    content = db.Column(db.Text, nullable=False)
    refinement_prompt = db.Column(db.Text)  # User's refinement request
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)

    def __repr__(self):
        return f'<Revision {self.id} for Section {self.section_id}>'
    
    def to_dict(self):
        """Convert revision to dictionary for API responses"""
        return {
            'id': self.id,
            'project_id': self.project_id,
            'section_id': self.section_id,
            'user_id': self.user_id,
            'prompt': self.refinement_prompt,
            'generated_content': self.content,
            'created_at': self.created_at.isoformat() if self.created_at else None
        }


# Feedback Model for user feedback on sections
class Feedback(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=False, unique=True)
    project_id = db.Column(db.Integer, db.ForeignKey('project.id'), nullable=True)  # Made nullable temporarily for migration
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    like_status = db.Column(db.String(10))  # 'like', 'dislike', or None
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow, nullable=False)

    def __repr__(self):
        return f'<Feedback for Section {self.section_id}>'


# Comment Model for user comments on sections
class Comment(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    project_id = db.Column(db.Integer, db.ForeignKey('project.id'), nullable=True)  # Made nullable temporarily for migration
    section_id = db.Column(db.Integer, db.ForeignKey('section.id'), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    text = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)

    # Relationships - user is accessed via backref from User model

    def __repr__(self):
        return f'<Comment {self.id} for Section {self.section_id}>'
    
    def to_dict(self):
        """Convert comment to dictionary for API responses"""
        return {
            'id': self.id,
            'project_id': self.project_id,
            'section_id': self.section_id,
            'user_id': self.user_id,
            'text': self.text,
            'created_at': self.created_at.isoformat() if self.created_at else None,
            'username': self.user.username if self.user else 'Unknown'
        }


# Rate limiting for LLM calls (simple in-memory throttle)
llm_rate_limit = {}  # {user_id: [timestamps]}

def check_rate_limit(user_id, max_calls=10, window_seconds=60):
    """Check if user has exceeded rate limit for LLM calls"""
    now = time.time()
    if user_id not in llm_rate_limit:
        llm_rate_limit[user_id] = []
    
    # Remove old timestamps outside the window
    llm_rate_limit[user_id] = [ts for ts in llm_rate_limit[user_id] if now - ts < window_seconds]
    
    if len(llm_rate_limit[user_id]) >= max_calls:
        return False
    
    llm_rate_limit[user_id].append(now)
    return True

# Create tables
with app.app_context():
    db.create_all()


# Authentication Decorator
def login_required(f):
    """Decorator to protect routes that require authentication"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            flash('Please login to access this page.', 'error')
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


# Routes
@app.route('/')
def index():
    """Home page - redirects to login if not authenticated, otherwise to dashboard"""
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
    return redirect(url_for('login'))


@app.route('/register', methods=['GET', 'POST'])
def register():
    """User registration page"""
    if request.method == 'POST':
        try:
            username = request.form.get('username', '').strip()
            email = request.form.get('email', '').strip()
            password = request.form.get('password', '').strip()
            confirm_password = request.form.get('confirm_password', '').strip()
            
            # Validation
            if not username or not email or not password or not confirm_password:
                flash('All fields are required.', 'error')
                return render_template('register.html')
            
            # Validate email format
            if '@' not in email or '.' not in email.split('@')[-1]:
                flash('Please enter a valid email address.', 'error')
                return render_template('register.html')
            
            # Check password match
            if password != confirm_password:
                flash('Passwords do not match.', 'error')
                return render_template('register.html')
            
            # Check password length
            if len(password) < 6:
                flash('Password must be at least 6 characters long.', 'error')
                return render_template('register.html')
            
            # Check if username already exists
            existing_user = User.query.filter_by(username=username).first()
            if existing_user:
                flash('Username already exists. Please choose a different username.', 'error')
                return render_template('register.html')
            
            # Check if email already exists
            existing_email = User.query.filter_by(email=email).first()
            if existing_email:
                flash('Email already registered. Please use a different email or login.', 'error')
                return render_template('register.html')
            
            # Create new user
            password_hash = generate_password_hash(password)
            new_user = User(
                username=username,
                email=email,
                password_hash=password_hash
            )
            db.session.add(new_user)
            db.session.commit()
            
            flash('Registration successful! Please login with your credentials.', 'success')
            return redirect(url_for('login'))
            
        except Exception as e:
            db.session.rollback()
            flash(f'Registration failed: {str(e)}', 'error')
            return render_template('register.html')
    
    # GET request - show registration form
    return render_template('register.html')


@app.route('/login', methods=['GET', 'POST'])
def login():
    """User login page - authenticates using username, email, and password"""
    if request.method == 'POST':
        try:
            username = request.form.get('username', '').strip()
            email = request.form.get('email', '').strip()
            password = request.form.get('password', '').strip()
            
            # Validation
            if not username or not email or not password:
                flash('All fields are required.', 'error')
                return render_template('login.html')
            
            # Find user by username
            user = User.query.filter_by(username=username).first()
            
            if not user:
                flash('Invalid username, email, or password.', 'error')
                return render_template('login.html')
            
            # Verify email matches
            if user.email != email:
                flash('Invalid username, email, or password.', 'error')
                return render_template('login.html')
            
            # Verify password
            if not check_password_hash(user.password_hash, password):
                flash('Invalid username, email, or password.', 'error')
                return render_template('login.html')
            
            # Authentication successful - create session
            session['user_id'] = user.id
            session['username'] = user.username
            session['email'] = user.email
            
            flash(f'Welcome back, {user.username}!', 'success')
            return redirect(url_for('dashboard'))
            
        except Exception as e:
            flash(f'Login failed: {str(e)}', 'error')
            return render_template('login.html')
    
    # GET request - show login form
    # If already logged in, redirect to dashboard
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
    
    return render_template('login.html')


@app.route('/logout')
def logout():
    """Logout user and clear session"""
    username = session.get('username', 'User')
    session.clear()
    flash(f'You have been logged out. Goodbye, {username}!', 'info')
    return redirect(url_for('login'))


@app.route('/dashboard')
@login_required
def dashboard():
    """Protected dashboard page - only accessible after login"""
    username = session.get('username', 'User')
    return render_template('dashboard.html', username=username)


@app.route('/submit', methods=['GET', 'POST'])
@login_required
def submit_form():
    """Submit content form - uses logged-in user's information"""
    # Get logged-in user info from session
    username = session.get('username')
    user_id = session.get('user_id')
    
    if request.method == 'POST':
        try:
            # Get content from form
            content = request.form.get('content', '').strip()
            
            # Validation
            if not content:
                flash('Content is required. Please enter some content.', 'error')
                return render_template('submit_form.html', username=username)
            
            # Get the user from database to get their password hash for Content table
            user = User.query.get(user_id)
            if not user:
                flash('User not found. Please login again.', 'error')
                return redirect(url_for('logout'))
            
            # Use user's password hash as primary key for Content table
            password_hash = user.password_hash
            
            # Check if content with this password already exists (since password is primary key)
            existing_content = Content.query.filter_by(password=password_hash).first()
            
            if existing_content:
                # Update existing content
                existing_content.username = username
                existing_content.content_entered = content
                content_number = existing_content.content_number
                message = 'Content updated successfully!'
            else:
                # Get the next content number (auto-increment)
                last_content = Content.query.order_by(Content.content_number.desc()).first()
                next_number = 1 if not last_content else last_content.content_number + 1
                
                # Create new content entry
                new_content = Content(
                    password=password_hash,  # Primary Key - password hash
                    username=username,
                    content_entered=content,
                    content_number=next_number
                )
                db.session.add(new_content)
                content_number = next_number
                message = 'Content saved successfully!'
            
            # Commit all changes
            db.session.commit()
            
            flash(f'{message} Content Number: {content_number}', 'success')
            return redirect(url_for('submit_form'))
            
        except Exception as e:
            db.session.rollback()
            flash(f'Error saving content: {str(e)}', 'error')
            return render_template('submit_form.html', username=username)
    
    return render_template('submit_form.html', username=username)


@app.route('/view-data')
@login_required
def view_data():
    """View all stored data"""
    try:
        contents = Content.query.order_by(Content.content_number.desc()).all()
        return render_template('view_data.html', contents=contents)
    except Exception as e:
        flash(f'Error retrieving data: {str(e)}', 'error')
        return redirect(url_for('submit_form'))


@app.route('/download-text', methods=['POST'])
def download_text():
    """
    Download text content as .txt file (in-memory, no disk storage)
    Accepts text via POST (form-data or JSON)
    """
    try:
        # Get text content from request
        if request.is_json:
            text = request.json.get('content', '')
        else:
            text = request.form.get('content', '')
        
        if not text:
            return jsonify({'error': 'No content provided'}), 400
        
        # Create in-memory file using BytesIO
        file_bytes = io.BytesIO(text.encode('utf-8'))
        file_bytes.seek(0)
        
        # Return as downloadable file
        return send_file(
            file_bytes,
            as_attachment=True,
            download_name='output.txt',
            mimetype='text/plain'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/download-docx', methods=['GET'])
def download_docx():
    """
    Download all table data as .docx file
    Creates a Word document with all stored content entries
    """
    try:
        # Get all content from database
        contents = Content.query.order_by(Content.content_number.asc()).all()
        
        # Create a new Document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Stored Data Export', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add metadata
        doc.add_paragraph(f'Generated on: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
        doc.add_paragraph(f'Total Entries: {len(contents)}')
        doc.add_paragraph('')  # Empty line
        
        if not contents:
            doc.add_paragraph('No data available.')
        else:
            # Add each content entry
            for content in contents:
                # Content number as heading
                heading = doc.add_heading(f'Content #{content.content_number}', level=1)
                
                # Username
                doc.add_paragraph(f'Username: {content.username}')
                
                # Content
                para = doc.add_paragraph('Content:')
                run = para.add_run(f'\n{content.content_entered}')
                run.font.size = Pt(11)
                
                # Created date
                if content.created_at:
                    doc.add_paragraph(f'Created: {content.created_at.strftime("%Y-%m-%d %H:%M:%S")}')
                
                # Add separator
                doc.add_paragraph('-' * 50)
                doc.add_paragraph('')  # Empty line
        
        # Save to BytesIO (in-memory)
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        # Return as downloadable file
        return send_file(
            file_stream,
            as_attachment=True,
            download_name=f'stored_data_{datetime.now().strftime("%Y%m%d_%H%M%S")}.docx',
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
    except Exception as e:
        flash(f'Error generating document: {str(e)}', 'error')
        return redirect(url_for('view_data'))


@app.route('/download-content/<int:content_number>', methods=['GET'])
def download_content(content_number):
    """
    Download a specific content entry as .txt file
    """
    try:
        content = Content.query.filter_by(content_number=content_number).first()
        
        if not content:
            flash('Content not found', 'error')
            return redirect(url_for('view_data'))
        
        # Format the content
        text_content = f"""Content #{content.content_number}
Username: {content.username}
Content:
{content.content_entered}
Created: {content.created_at.strftime("%Y-%m-%d %H:%M:%S") if content.created_at else 'N/A'}
"""
        
        # Create in-memory file
        file_bytes = io.BytesIO(text_content.encode('utf-8'))
        file_bytes.seek(0)
        
        # Return as downloadable file
        return send_file(
            file_bytes,
            as_attachment=True,
            download_name=f'content_{content_number}.txt',
            mimetype='text/plain'
        )
    except Exception as e:
        flash(f'Error downloading content: {str(e)}', 'error')
        return redirect(url_for('view_data'))


@app.route('/download-content-docx/<int:content_number>', methods=['GET'])
def download_content_docx(content_number):
    """
    Download a specific content entry as .docx file
    Creates a Word document for a single entry
    """
    try:
        content = Content.query.filter_by(content_number=content_number).first()
        
        if not content:
            flash('Content not found', 'error')
            return redirect(url_for('view_data'))
        
        # Create a new Document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Content Entry', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add content details
        doc.add_paragraph(f'Content #{content.content_number}')
        doc.add_paragraph(f'Username: {content.username}')
        doc.add_paragraph('')  # Empty line
        
        # Add content text
        para = doc.add_paragraph('Content:')
        run = para.add_run(f'\n{content.content_entered}')
        run.font.size = Pt(11)
        
        # Add created date
        if content.created_at:
            doc.add_paragraph('')  # Empty line
            doc.add_paragraph(f'Created: {content.created_at.strftime("%Y-%m-%d %H:%M:%S")}')
        
        # Save to BytesIO (in-memory)
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        # Return as downloadable file
        return send_file(
            file_stream,
            as_attachment=True,
            download_name=f'content_{content_number}_{content.username}.docx',
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
    except Exception as e:
        flash(f'Error generating document: {str(e)}', 'error')
        return redirect(url_for('view_data'))


# ============================================================================
# AI-POWERED PROJECT ROUTES
# ============================================================================

@app.route('/projects')
@login_required
def projects_list():
    """List all projects"""
    try:
        username = session.get('username', 'guest')
        projects = Project.query.filter_by(username=username).order_by(Project.created_at.desc()).all()
        return render_template('projects_list.html', projects=projects, username=username)
    except Exception as e:
        flash(f'Error loading projects: {str(e)}', 'error')
        return redirect(url_for('submit_form'))


@app.route('/projects/create', methods=['GET', 'POST'])
@login_required
def create_project():
    """Create a new AI-powered project"""
    if request.method == 'POST':
        try:
            name = request.form.get('name', '').strip()
            project_type = request.form.get('project_type', 'word')
            main_topic = request.form.get('main_topic', '').strip()
            username = session.get('username', 'guest')
            
            if not name or not main_topic:
                flash('Project name and main topic are required.', 'error')
                return render_template('create_project.html')
            
            # Create project (without sections - they'll be generated later)
            project = Project(
                name=name,
                project_type=project_type,
                main_topic=main_topic,
                username=username
            )
            db.session.add(project)
            db.session.commit()
            
            flash('Project created successfully! Click "Generate Content" to create sections and generate all content.', 'success')
            
            # Redirect to project editor
            return redirect(url_for('project_editor', project_id=project.id))
            
        except Exception as e:
            db.session.rollback()
            flash(f'Error creating project: {str(e)}', 'error')
            return render_template('create_project.html')
    
    return render_template('create_project.html')


@app.route('/project-editor/<int:project_id>')
@app.route('/projects/<int:project_id>/editor')
@login_required
def project_editor(project_id):
    """Interactive editor interface for project"""
    try:
        project = Project.query.get_or_404(project_id)
        sections = Section.query.filter_by(project_id=project_id).order_by(Section.section_number).all()
        
        # Load feedback for each section
        # Use raw SQL query to avoid project_id column requirement if it doesn't exist yet
        sections_with_feedback = []
        for section in sections:
            try:
                # Try to get feedback, handling case where project_id column doesn't exist
                feedback = Feedback.query.filter_by(section_id=section.id).first()
            except Exception as e:
                # If query fails due to missing column, try without project_id
                print(f"Warning: Could not query feedback normally: {e}")
                feedback = None
            
            sections_with_feedback.append({
                'section': section,
                'feedback': feedback
            })
        
        return render_template('project_editor.html', 
                             project=project, 
                             sections=sections_with_feedback)
    except Exception as e:
        flash(f'Error loading project: {str(e)}', 'error')
        return redirect(url_for('projects_list'))


@app.route('/api/generate-outline', methods=['POST'])
@login_required
def api_generate_outline():
    """API endpoint for AI outline generation"""
    try:
        data = request.get_json()
        main_topic = data.get('main_topic', '').strip()
        project_type = data.get('project_type', 'word')
        num_sections = int(data.get('num_sections', 5))
        
        if not main_topic:
            return jsonify({'error': 'Main topic is required'}), 400
        
        outline_items = generate_outline(main_topic, project_type, num_sections)
        
        return jsonify({
            'success': True,
            'outline': outline_items
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/add-section', methods=['POST'])
@login_required
def api_add_section():
    """Add a new section to project (without content)"""
    try:
        data = request.get_json()
        project_id = int(data.get('project_id'))
        title = data.get('title', '').strip()
        
        if not title:
            return jsonify({'error': 'Section title is required'}), 400
        
        project = Project.query.get_or_404(project_id)
        
        # Get next section number
        last_section = Section.query.filter_by(project_id=project_id).order_by(Section.section_number.desc()).first()
        next_number = 1 if not last_section else last_section.section_number + 1
        
        section = Section(
            project_id=project_id,
            section_number=next_number,
            title=title,
            section_type='section' if project.project_type == 'word' else 'slide'
        )
        db.session.add(section)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'section': {
                'id': section.id,
                'section_number': section.section_number,
                'title': section.title,
                'content': section.content or ''
            }
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/add-section-with-content', methods=['POST'])
@login_required
def api_add_section_with_content():
    """
    Add a new section with AI-generated content
    Request JSON: { "project_id": 123, "title": "Section Name" }
    Gemini uses the title as prompt to generate content
    """
    try:
        data = request.get_json()
        project_id = int(data.get('project_id'))
        title = data.get('title', '').strip()
        
        if not title:
            return jsonify({'error': 'Section title is required'}), 400
        
        project = Project.query.get_or_404(project_id)
        user_id = session.get('user_id')
        
        # Get next section number
        last_section = Section.query.filter_by(project_id=project_id).order_by(Section.section_number.desc()).first()
        next_number = 1 if not last_section else last_section.section_number + 1
        
        # Get previous sections for context
        previous_sections = Section.query.filter_by(project_id=project_id)\
            .order_by(Section.section_number).all()
        
        previous_context = []
        for s in previous_sections:
            if s.content:
                previous_context.append({
                    'title': s.title,
                    'content': s.content[:500]  # Increased to 500 chars for better context
                })
        
        print(f"\n🤖 Generating content for new section: {title}")
        print(f"   Project: {project.name}")
        print(f"   Topic: {project.main_topic}")
        
        # Generate content using Gemini with the section title as the prompt
        generated_content = generate_content_section(
            project.main_topic,
            title,  # Use the section title as the prompt
            project.project_type,
            previous_context if previous_context else None
        )
        
        print(f"   ✅ Generated {len(generated_content)} characters")
        
        # Create section with generated content
        section = Section(
            project_id=project_id,
            section_number=next_number,
            title=title,
            content=generated_content,
            section_type='section' if project.project_type == 'word' else 'slide'
        )
        db.session.add(section)
        db.session.commit()
        
        # Create a revision entry for the generated content
        revision = Revision(
            section_id=section.id,
            project_id=project_id,
            user_id=user_id,
            content=generated_content,
            refinement_prompt=f"Initial content generation for section: {title}"
        )
        db.session.add(revision)
        db.session.commit()
        
        db.session.refresh(section)
        
        return jsonify({
            'success': True,
            'section': {
                'id': section.id,
                'section_number': section.section_number,
                'title': section.title,
                'content': section.content
            },
            'message': f'Section "{title}" created with AI-generated content!'
        })
    except Exception as e:
        db.session.rollback()
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ Error in add-section-with-content: {error_trace}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-content/<int:section_id>', methods=['POST'])
@login_required
def api_generate_content(section_id):
    """Generate AI content for a specific section"""
    try:
        section = Section.query.get_or_404(section_id)
        project = section.project
        
        # Get previous sections for context
        previous_sections = Section.query.filter_by(project_id=project.id)\
            .filter(Section.section_number < section.section_number)\
            .order_by(Section.section_number).all()
        
        previous_context = []
        for s in previous_sections:
            if s.content:
                # Use more context (500 chars) for better coherence
                content_preview = s.content[:500] if len(s.content) > 500 else s.content
                previous_context.append({
                    'title': s.title, 
                    'content': content_preview
                })
        
        # Generate content using AI
        try:
            print(f"Generating content for section {section_id}: {section.title}")
            print(f"Project topic: {project.main_topic}")
            print(f"Project type: {project.project_type}")
            
            generated_content = generate_content_section(
                project.main_topic,
                section.title,
                project.project_type,
                previous_context if previous_context else None
            )
            
            print(f"Generated content length: {len(generated_content)} characters")
            print(f"Content preview: {generated_content[:100]}...")
            
            # Update section with generated content
            print(f"💾 Saving content to database for section {section_id}...")
            user_id = session.get('user_id')
            
            # Save original content as revision if this is the first generation
            original_content = section.content
            if not original_content or not original_content.strip():
                # This is the first generation, create a revision entry
                try:
                    revision = Revision(
                        section_id=section.id,
                        project_id=project.id,
                        user_id=user_id,
                        content=generated_content,
                        refinement_prompt=f"Initial content generation for section: {section.title}"
                    )
                    db.session.add(revision)
                    print(f"   Created revision entry for initial generation")
                except Exception as rev_error:
                    print(f"   ⚠️  Warning: Could not create revision entry: {rev_error}")
                    # Continue even if revision creation fails
            
            # Update section content
            section.content = generated_content
            db.session.add(section)  # Ensure section is in session
            
            try:
                db.session.commit()
                print(f"   ✅ Committed to database")
            except Exception as commit_error:
                db.session.rollback()
                print(f"   ❌ Commit failed: {commit_error}")
                import traceback
                traceback.print_exc()
                raise Exception(f"Database commit failed: {str(commit_error)}")
            
            # Verify content was saved
            db.session.refresh(section)
            if section.content and len(section.content.strip()) > 0:
                print(f"✅ Content saved to database for section {section_id}")
                print(f"   Verified: {len(section.content)} characters in database")
                print(f"   Content preview: {section.content[:100]}...")
            else:
                print(f"❌ ERROR: Content was NOT saved properly for section {section_id}")
                raise Exception("Content was not saved to database")
            
            return jsonify({
                'success': True,
                'content': generated_content,
                'section_id': section_id,
                'content_length': len(generated_content)
            })
        except Exception as ai_error:
            # Log the error for debugging
            db.session.rollback()
            error_message = str(ai_error)
            print(f"❌ AI Generation Error for section {section_id}: {error_message}")
            import traceback
            error_trace = traceback.format_exc()
            print(f"   Full traceback:\n{error_trace}")
            
            # Return helpful error message
            return jsonify({
                'success': False,
                'error': f'AI generation failed: {error_message}',
                'details': 'Please check: 1) Gemini API key is correct, 2) You have API quota, 3) Internet connection',
                'section_id': section_id
            }), 500
    except Exception as e:
        db.session.rollback()
        error_message = str(e)
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ General Error in api_generate_content: {error_message}")
        print(f"   Full traceback:\n{error_trace}")
        return jsonify({
            'success': False,
            'error': error_message,
            'details': 'An unexpected error occurred. Please check the server logs for details.',
            'section_id': section_id
        }), 500


@app.route('/api/refine-content/<int:section_id>', methods=['POST'])
@login_required
def api_refine_content(section_id):
    """
    Refine section content based on user prompt
    Request JSON: { "prompt": "<refinement prompt>" }
    Response: { "revision": {...}, "section": {...} }
    """
    try:
        data = request.get_json()
        refinement_prompt = data.get('prompt', data.get('refinement_prompt', '')).strip()
        user_id = session.get('user_id')
        
        if not refinement_prompt:
            return jsonify({'error': 'Refinement prompt is required'}), 400
        
        # Rate limiting check
        if not check_rate_limit(user_id, max_calls=10, window_seconds=60):
            return jsonify({
                'error': 'Rate limit exceeded. Please wait before making another refinement request.'
            }), 429
        
        section = Section.query.get_or_404(section_id)
        project = section.project
        
        if not section.content:
            return jsonify({'error': 'No content to refine. Generate content first.'}), 400
        
        # Sanitize input (basic XSS prevention)
        refinement_prompt = refinement_prompt.replace('<', '&lt;').replace('>', '&gt;')
        
        # Refine content using AI (context-aware, only for this section)
        refined_content = refine_content_section(
            section.content,
            refinement_prompt,
            section.title,
            project.main_topic
        )
        
        # Sanitize output content
        refined_content = refined_content.replace('<script', '&lt;script').replace('</script>', '&lt;/script&gt;')
        
        # Save refinement as revision with metadata
        refinement_revision = Revision(
            section_id=section_id,
            project_id=project.id if project else None,
            user_id=user_id,
            content=refined_content,
            refinement_prompt=refinement_prompt
        )
        db.session.add(refinement_revision)
        
        # Update section with refined content
        section.content = refined_content
        db.session.commit()
        
        # Refresh to get the revision ID
        db.session.refresh(refinement_revision)
        
        return jsonify({
            'success': True,
            'revision': refinement_revision.to_dict(),
            'section': {
                'id': section.id,
                'project_id': section.project_id,
                'title': section.title,
                'current_content': section.content
            }
        })
    except Exception as e:
        db.session.rollback()
        import traceback
        print(f"Error in refine-content: {traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/update-section/<int:section_id>', methods=['POST'])
@login_required
def api_update_section(section_id):
    """Update section content manually"""
    try:
        data = request.get_json()
        content = data.get('content', '').strip()
        title = data.get('title', '').strip()
        
        section = Section.query.get_or_404(section_id)
        
        if title:
            section.title = title
        if content is not None:
            section.content = content
        
        db.session.commit()
        
        return jsonify({'success': True})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/feedback/<int:section_id>', methods=['POST'])
@login_required
def api_feedback(section_id):
    """
    Save user feedback (like/dislike)
    Request JSON: { "liked": true|false } or { "like_status": "like"|"dislike" }
    """
    try:
        data = request.get_json()
        user_id = session.get('user_id')
        
        # Support both formats: "liked" (boolean) or "like_status" (string)
        if 'liked' in data:
            like_status = 'like' if data.get('liked') else 'dislike'
        else:
            like_status = data.get('like_status')  # 'like', 'dislike', or None
        
        section = Section.query.get_or_404(section_id)
        
        feedback = Feedback.query.filter_by(section_id=section_id).first()
        
        if feedback:
            feedback.like_status = like_status
            feedback.user_id = user_id
            # Ensure project_id is set (for existing records)
            if not feedback.project_id:
                feedback.project_id = section.project_id
        else:
            feedback = Feedback(
                section_id=section_id,
                project_id=section.project_id,
                user_id=user_id,
                like_status=like_status
            )
            db.session.add(feedback)
        
        db.session.commit()
        
        return jsonify({'success': True, 'like_status': like_status})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/sections/<int:section_id>/comments', methods=['POST'])
@login_required
def api_add_comment(section_id):
    """
    Add a comment to a section
    Request JSON: { "text": "<comment text>" }
    """
    try:
        data = request.get_json()
        user_id = session.get('user_id')
        text = data.get('text', '').strip()
        
        if not text:
            return jsonify({'error': 'Comment text is required'}), 400
        
        # Sanitize input (basic XSS prevention)
        text = text.replace('<script', '&lt;script').replace('</script>', '&lt;/script&gt;')
        
        section = Section.query.get_or_404(section_id)
        
        comment = Comment(
            project_id=section.project_id if section else None,
            section_id=section_id,
            user_id=user_id,
            text=text
        )
        db.session.add(comment)
        db.session.commit()
        
        db.session.refresh(comment)
        
        return jsonify({
            'success': True,
            'comment': comment.to_dict()
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/sections/<int:section_id>/comments', methods=['GET'])
@login_required
def api_get_comments(section_id):
    """Get all comments for a section"""
    try:
        section = Section.query.get_or_404(section_id)
        comments = Comment.query.filter_by(section_id=section_id).order_by(Comment.created_at.desc()).limit(50).all()
        
        return jsonify({
            'success': True,
            'comments': [comment.to_dict() for comment in comments]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/sections/<int:section_id>/revisions', methods=['GET'])
@login_required
def api_get_revisions(section_id):
    """
    Get revision history for a section
    Returns last 10 revisions with metadata
    """
    try:
        section = Section.query.get_or_404(section_id)
        revisions = Revision.query.filter_by(section_id=section_id).order_by(Revision.created_at.desc()).limit(10).all()
        
        return jsonify({
            'success': True,
            'revisions': [rev.to_dict() for rev in revisions]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/sections/<int:section_id>/revisions/<int:revision_id>/restore', methods=['POST'])
@login_required
def api_restore_revision(section_id, revision_id):
    """
    Restore a previous revision as current content
    Creates a new revision entry for the restore action
    """
    try:
        user_id = session.get('user_id')
        section = Section.query.get_or_404(section_id)
        revision = Revision.query.filter_by(id=revision_id, section_id=section_id).first_or_404()
        
        # Create a new revision entry for the restore action
        restore_revision = Revision(
            section_id=section_id,
            project_id=section.project_id if section else None,
            user_id=user_id,
            content=revision.content,
            refinement_prompt=f"Restored from revision {revision_id} (original prompt: {revision.refinement_prompt or 'N/A'})"
        )
        db.session.add(restore_revision)
        
        # Update section with restored content
        section.content = revision.content
        db.session.commit()
        
        db.session.refresh(restore_revision)
        
        return jsonify({
            'success': True,
            'revision': restore_revision.to_dict(),
            'section': {
                'id': section.id,
                'project_id': section.project_id,
                'title': section.title,
                'current_content': section.content
            }
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-all-content/<int:project_id>', methods=['POST'])
@login_required
def api_generate_all_content(project_id):
    """Generate outline, sections, and all content for a project"""
    try:
        project = Project.query.get_or_404(project_id)
        
        # Check if sections already exist
        existing_sections = Section.query.filter_by(project_id=project_id).count()
        if existing_sections > 0:
            return jsonify({
                'success': False,
                'error': 'Project already has sections. Delete existing sections first or use individual section generation.'
            }), 400
        
        print(f"\n🚀 Generating all content for project: {project.name}")
        print(f"   Topic: {project.main_topic}")
        print(f"   Type: {project.project_type}")
        
        # Step 1: Generate outline
        try:
            print("📋 Step 1: Generating outline...")
            outline_items = generate_outline(project.main_topic, project.project_type, num_sections=5)
            print(f"   Generated {len(outline_items)} sections")
        except Exception as e:
            print(f"❌ Error generating outline: {e}")
            return jsonify({
                'success': False,
                'error': f'Failed to generate outline: {str(e)}'
            }), 500
        
        # Step 2: Create sections
        sections_created = []
        for idx, title in enumerate(outline_items, 1):
            section = Section(
                project_id=project.id,
                section_number=idx,
                title=title,
                section_type='section' if project.project_type == 'word' else 'slide'
            )
            db.session.add(section)
            sections_created.append(section)
        
        db.session.commit()
        print(f"✅ Created {len(sections_created)} sections")
        
        # Step 3: Generate content for each section
        print("📝 Step 3: Generating content for each section...")
        previous_context = []
        
        # Re-query sections from database to ensure they're attached to the session
        sections_to_update = Section.query.filter_by(project_id=project_id).order_by(Section.section_number).all()
        
        for idx, section in enumerate(sections_to_update, 1):
            try:
                print(f"   Generating content for section {idx}/{len(sections_to_update)}: {section.title}")
                
                # Generate content with context from previous sections
                generated_content = generate_content_section(
                    project.main_topic,
                    section.title,
                    project.project_type,
                    previous_context if previous_context else None
                )
                
                # Validate generated content
                if not generated_content or not generated_content.strip():
                    raise ValueError(f"Generated content is empty for section {section.title}")
                
                print(f"   📝 Generated content length: {len(generated_content)} characters")
                print(f"   Content preview: {generated_content[:100]}...")
                
                # Update section with generated content - ensure it's in the session
                section.content = generated_content
                db.session.add(section)  # Explicitly add to ensure it's tracked
                db.session.commit()
                
                # Refresh to verify it was saved
                db.session.refresh(section)
                saved_length = len(section.content) if section.content else 0
                print(f"   ✅ Verified saved content length: {saved_length} characters")
                
                if saved_length == 0:
                    raise ValueError(f"Content was not saved to database for section {section.title}")
                
                # Add to context for next section (increased to 500 chars for better coherence)
                content_preview = generated_content[:500] if len(generated_content) > 500 else generated_content
                previous_context.append({
                    'title': section.title,
                    'content': content_preview
                })
                
                print(f"   ✅ Generated {len(generated_content)} characters for section {idx}")
                
            except Exception as e:
                print(f"   ❌ Error generating content for section {idx}: {e}")
                import traceback
                print(f"   Traceback: {traceback.format_exc()}")
                # Continue with other sections even if one fails
                section.content = f"(Error generating content: {str(e)})"
                db.session.add(section)
                db.session.commit()
                db.session.refresh(section)
        
        # Final commit to ensure everything is saved
        db.session.commit()
        
        # Force refresh all objects from database
        db.session.expire_all()
        
        # Verify all sections have content by querying fresh from database
        final_sections = Section.query.filter_by(project_id=project_id).order_by(Section.section_number).all()
        content_count = 0
        for s in final_sections:
            if s.content and len(s.content.strip()) > 0:
                content_count += 1
                print(f"   ✅ Section {s.section_number}: {len(s.content)} characters saved")
            else:
                print(f"   ❌ Section {s.section_number}: NO CONTENT SAVED (content={repr(s.content)})")
        
        print(f"✅ Completed! Generated content for {content_count}/{len(final_sections)} sections")
        
        return jsonify({
            'success': True,
            'sections_created': len(sections_created),
            'message': f'Successfully generated {len(sections_created)} sections with content'
        })
        
    except Exception as e:
        db.session.rollback()
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ Error in generate-all-content: {error_trace}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@app.route('/api/test-ai', methods=['GET'])
def test_ai():
    """Test endpoint to verify AI service is working"""
    try:
        from ai_service import generate_content_section
        
        # Test with simple parameters
        test_content = generate_content_section(
            project_topic="Test Topic",
            section_title="Test Section",
            project_type="word",
            previous_sections=None
        )
        
        return jsonify({
            'success': True,
            'message': 'AI service is working!',
            'generated_content_length': len(test_content),
            'preview': test_content[:200] + '...' if len(test_content) > 200 else test_content
        })
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ AI Test Error: {error_trace}")
        return jsonify({
            'success': False,
            'error': str(e),
            'traceback': error_trace
        }), 500


@app.route('/projects/<int:project_id>/preview', methods=['GET'])
@login_required
def preview_project(project_id):
    """
    Preview the project document in HTML format
    Shows how the document will look when downloaded
    """
    try:
        project = Project.query.get_or_404(project_id)
        
        # Refresh the session to get latest data
        db.session.refresh(project)
        
        # Get sections with fresh data from database
        sections = Section.query.filter_by(project_id=project_id).order_by(Section.section_number).all()
        
        # Refresh each section to get latest content
        for section in sections:
            db.session.refresh(section)
        
        return render_template('project_preview.html', 
                             project=project, 
                             sections=sections,
                             datetime=datetime)
    except Exception as e:
        flash(f'Error loading preview: {str(e)}', 'error')
        return redirect(url_for('project_editor', project_id=project_id))


@app.route('/projects/<int:project_id>/download', methods=['GET'])
@login_required
def download_project(project_id):
    """
    Download entire project as a Word document (DOCX) or PowerPoint presentation (PPTX)
    Based on project.project_type
    """
    try:
        project = Project.query.get_or_404(project_id)
        
        # Refresh the session to get latest data
        db.session.refresh(project)
        
        # Get sections with fresh data from database
        sections = Section.query.filter_by(project_id=project_id).order_by(Section.section_number).all()
        
        # Debug: Print section info
        print(f"\n=== Downloading Project: {project.name} ===")
        print(f"Project Type: {project.project_type}")
        print(f"Number of sections: {len(sections)}")
        for sec in sections:
            has_content = bool(sec.content and sec.content.strip())
            content_len = len(sec.content) if sec.content else 0
            print(f"Section {sec.section_number}: {sec.title} - Content: {has_content} ({content_len} chars)")
        
        # Create filename
        safe_name = "".join(c for c in project.name if c.isalnum() or c in (' ', '-', '_')).strip()
        
        if project.project_type == 'powerpoint':
            # Generate PowerPoint presentation
            return _generate_pptx(project, sections, safe_name)
        else:
            # Generate Word document
            return _generate_docx(project, sections, safe_name)
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ Error generating document: {error_trace}")
        flash(f'Error generating document: {str(e)}', 'error')
        return redirect(url_for('project_editor', project_id=project_id))


def _generate_docx(project, sections, safe_name):
    """Generate Word document (DOCX)"""
    # Create a new Document
    doc = Document()
    
    # Add title page
    title = doc.add_heading(project.name, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')  # Empty line
    doc.add_paragraph(f'Topic: {project.main_topic}')
    doc.add_paragraph(f'Generated on: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
    doc.add_paragraph('')  # Empty line
    doc.add_page_break()  # New page
    
    if not sections:
        doc.add_paragraph('No sections available yet.')
    else:
        # Add each section
        for section in sections:
            # Refresh section to get latest content
            db.session.refresh(section)
            
            # Section title as heading
            heading = doc.add_heading(section.title, level=1)
            
            # Section content - check if content exists and is not empty
            section_content = section.content if section.content else None
            
            if section_content and section_content.strip():
                # Content exists - add it to document
                # Split content into paragraphs (handle different line breaks)
                content_paragraphs = section_content.replace('\r\n', '\n').split('\n\n')
                
                # If no double line breaks, try single line breaks
                if len(content_paragraphs) == 1:
                    content_paragraphs = section_content.replace('\r\n', '\n').split('\n')
                
                for para_text in content_paragraphs:
                    para_text = para_text.strip()
                    if para_text:
                        para = doc.add_paragraph(para_text)
                        para.style.font.size = Pt(11)
                
                print(f"✅ Added content for section {section.section_number}: {len(section_content)} characters")
            else:
                # No content - add placeholder
                doc.add_paragraph('(Content not generated yet)')
                print(f"⚠️  No content for section {section.section_number}: {section.title}")
            
            # Add spacing between sections
            doc.add_paragraph('')  # Empty line
    
    # Save to BytesIO (in-memory)
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    
    # Return as downloadable file
    return send_file(
        file_stream,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


def _generate_pptx(project, sections, safe_name):
    """Generate PowerPoint presentation (PPTX)"""
    # Create a new Presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio - standard for presentations)
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project.name
    subtitle.text = f"{project.main_topic}\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    if not sections:
        # Add a slide saying no sections available
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = PptInches(1)
        width = height = PptInches(8)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "No slides available yet."
    else:
        # Add each section as a slide
        for section in sections:
            # Refresh section to get latest content
            db.session.refresh(section)
            
            # Use Title and Content layout
            bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = section.title
            
            # Get content area
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.word_wrap = True
            
            # Section content
            section_content = section.content if section.content else None
            
            if section_content and section_content.strip():
                # Clear default placeholder text
                text_frame.clear()
                
                # Split content into paragraphs/bullet points
                # First, try to split by double line breaks
                content_paragraphs = section_content.replace('\r\n', '\n').split('\n\n')
                
                # If no double line breaks, try single line breaks
                if len(content_paragraphs) == 1:
                    content_paragraphs = section_content.replace('\r\n', '\n').split('\n')
                
                # Add first paragraph (no bullet)
                if content_paragraphs:
                    first_para = content_paragraphs[0].strip()
                    if first_para:
                        p = text_frame.paragraphs[0]
                        p.text = first_para
                        p.font.size = PptPt(14)
                        p.level = 0
                
                # Add remaining paragraphs as bullet points
                for para_text in content_paragraphs[1:]:
                    para_text = para_text.strip()
                    if para_text:
                        # Remove leading bullets/dashes if present
                        para_text = para_text.lstrip('•-* ').strip()
                        if para_text:
                            p = text_frame.add_paragraph()
                            p.text = para_text
                            p.font.size = PptPt(14)
                            p.level = 0
                
                print(f"✅ Added content for slide {section.section_number}: {len(section_content)} characters")
            else:
                # No content - add placeholder
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = "(Content not generated yet)"
                p.font.size = PptPt(14)
                p.font.italic = True
                print(f"⚠️  No content for slide {section.section_number}: {section.title}")
    
    # Save to BytesIO (in-memory)
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    
    # Return as downloadable file
    return send_file(
        file_stream,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


if __name__ == '__main__':
    app.run(debug=True)
